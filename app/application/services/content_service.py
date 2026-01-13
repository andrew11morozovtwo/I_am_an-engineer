"""
Сервис для обработки контента из различных источников.

Извлекает текст из URL, обрабатывает изображения, PDF, аудио и другие типы медиа.
"""
import io
import logging
import re
import http.client
import asyncio
from urllib.parse import urlparse
from typing import Optional

import aiohttp
from bs4 import BeautifulSoup
from pypdf import PdfReader
from aiogram import Bot, types

# Импорты для обработки документов (опциональные, чтобы не падать если библиотеки не установлены)
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

logger = logging.getLogger(__name__)

# Предупреждения о недоступных библиотеках (после определения logger)
if DocxDocument is None:
    logger.warning("python-docx не установлен, обработка Word документов будет недоступна")
if load_workbook is None:
    logger.warning("openpyxl не установлен, обработка Excel файлов будет недоступна")
if Presentation is None:
    logger.warning("python-pptx не установлен, обработка PowerPoint файлов будет недоступна")


logger = logging.getLogger(__name__)


async def extract_text_from_url(url: str) -> Optional[str]:
    """
    Извлекает текст с веб-страницы по URL.

    :param url: URL веб-страницы
    :return: Извлеченный текст или None при ошибке
    """
    try:
        parsed_url = urlparse(url)
        if not parsed_url.netloc:
            logger.warning(f"Некорректный URL: {url}")
            return None

        async with aiohttp.ClientSession() as session:
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
            }
            async with session.get(url, headers=headers, timeout=aiohttp.ClientTimeout(total=10)) as response:
                if response.status == 200:
                    content = await response.text()
                    soup = BeautifulSoup(content, 'html.parser')
                    text_content = soup.get_text()
                    cleaned_text = "\n".join(
                        line.strip() for line in text_content.splitlines() if line.strip()
                    )
                    return cleaned_text
                else:
                    logger.warning(f"Ошибка при получении страницы {url}: статус {response.status}")
                    return None
    except Exception as e:
        logger.error(f"Ошибка при извлечении текста из URL {url}: {e}")
        return None


def find_url_in_text(text: str) -> Optional[str]:
    """
    Находит первую URL в тексте.

    :param text: Текст для поиска
    :return: Найденный URL или None
    """
    if not text:
        return None
    url_match = re.search(r'(http[s]?://[^\s]+)', text)
    return url_match.group(0) if url_match else None


async def process_url_in_text(text: str) -> str:
    """
    Обрабатывает текст, извлекая контент из найденных URL.

    :param text: Исходный текст
    :return: Объединенный текст (исходный + текст с веб-страницы)
    """
    url = find_url_in_text(text)
    if url:
        extracted_text = await extract_text_from_url(url)
        if extracted_text:
            return f"{text}\n\n{extracted_text}"
    return text


async def get_image_description(image_url: str, openai_client) -> Optional[str]:
    """
    Получает описание изображения через OpenAI Vision API.

    :param image_url: URL изображения
    :param openai_client: Клиент OpenAI
    :return: Описание изображения или None при ошибке
    """
    try:
        logger.info(f"Запрашиваем описание изображения через Vision API: {image_url[:100]}...")
        # Оборачиваем синхронный вызов OpenAI в asyncio.to_thread для неблокирующего выполнения
        response = await asyncio.to_thread(
            openai_client.chat.completions.create,
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Что на этом изображении? Дай краткое описание на русском языке."},
                        {
                            "type": "image_url",
                            "image_url": {"url": image_url},
                        },
                    ],
                }
            ],
            max_tokens=500,
        )
        description = response.choices[0].message.content
        logger.info(f"Описание изображения получено: {description[:100]}...")
        return description
    except Exception as e:
        logger.error(f"Ошибка при обращении к OpenAI Vision API: {e}", exc_info=True)
        return None


async def get_photo_url(bot: Bot, message: types.Message) -> Optional[str]:
    """
    Получает URL фотографии из сообщения Telegram.

    :param bot: Экземпляр бота
    :param message: Сообщение с фотографией
    :return: URL фотографии или None
    """
    try:
        if message.photo:
            file_id = message.photo[-1].file_id
            file_info = await bot.get_file(file_id)
            return f'https://api.telegram.org/file/bot{bot.token}/{file_info.file_path}'
    except Exception as e:
        logger.error(f"Ошибка при получении URL изображения: {e}")
    return None


async def extract_pdf_text(pdf_url: str) -> Optional[str]:
    """
    Извлекает текст из PDF файла по URL.

    :param pdf_url: URL PDF файла
    :return: Извлеченный текст или None при ошибке
    """
    logger.info(f"Начинаем извлечение текста из PDF по URL: {pdf_url[:100]}...")
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(pdf_url, timeout=aiohttp.ClientTimeout(total=60)) as response:
                response.raise_for_status()
                pdf_content = await response.read()
                logger.info(f"PDF файл скачан, размер: {len(pdf_content)} байт")
                pdf_file = io.BytesIO(pdf_content)
                pdf_reader = PdfReader(pdf_file)

                pdf_text = ""
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            pdf_text += f"\n--- Страница {page_num} ---\n{page_text}"
                    except Exception as e:
                        logger.warning(f"Не удалось извлечь текст со страницы {page_num}: {e}")
                        continue

                if not pdf_text.strip():
                    logger.warning(f"Из PDF {pdf_url[:100]}... не удалось извлечь текст.")
                    return None

                # Ограничиваем размер текста (примерно 4000 токенов)
                if len(pdf_text) > 12000:
                    pdf_text = pdf_text[:12000] + "\n... (текст обрезан из-за ограничений)"
                
                logger.info(f"Текст из PDF {pdf_url[:100]}... успешно извлечен ({len(pdf_text)} символов).")
                return pdf_text
    except asyncio.TimeoutError:
        logger.error(f"Таймаут при скачивании или обработке PDF по URL: {pdf_url[:100]}...")
        return None
    except Exception as e:
        logger.error(f"Ошибка при обработке PDF файла {pdf_url[:100]}...: {e}", exc_info=True)
        return None


async def extract_document_text(document_url: str, file_extension: str) -> Optional[str]:
    """
    Извлекает текст из различных типов документов по URL.

    Поддерживаемые форматы:
    - .txt - текстовые файлы
    - .docx - Word документы
    - .xlsx - Excel файлы
    - .pptx - PowerPoint презентации
    - .odt - OpenDocument Text

    :param document_url: URL документа
    :param file_extension: Расширение файла (например, '.docx', '.xlsx')
    :return: Извлеченный текст или None при ошибке
    """
    logger.info(f"Начинаем извлечение текста из документа {file_extension} по URL: {document_url[:100]}...")
    try:
        # Скачиваем файл
        async with aiohttp.ClientSession() as session:
            async with session.get(document_url, timeout=aiohttp.ClientTimeout(total=60)) as response:
                response.raise_for_status()
                document_content = await response.read()
                logger.info(f"Документ скачан, размер: {len(document_content)} байт")
                document_file = io.BytesIO(document_content)
                
                document_text = ""
                file_ext_lower = file_extension.lower()
                
                # Обработка текстовых файлов
                if file_ext_lower == '.txt':
                    try:
                        document_file.seek(0)
                        # Пробуем разные кодировки
                        for encoding in ['utf-8', 'cp1251', 'windows-1251', 'latin-1']:
                            try:
                                document_file.seek(0)
                                document_text = document_file.read().decode(encoding)
                                break
                            except UnicodeDecodeError:
                                continue
                        if not document_text:
                            logger.warning("Не удалось декодировать текстовый файл")
                            return None
                    except Exception as e:
                        logger.error(f"Ошибка при обработке текстового файла: {e}", exc_info=True)
                        return None
                
                # Обработка Word документов (.docx)
                elif file_ext_lower == '.docx':
                    try:
                        from docx import Document
                        document_file.seek(0)
                        doc = Document(document_file)
                        for paragraph in doc.paragraphs:
                            if paragraph.text.strip():
                                document_text += paragraph.text + "\n"
                        # Также извлекаем текст из таблиц
                        for table in doc.tables:
                            for row in table.rows:
                                for cell in row.cells:
                                    if cell.text.strip():
                                        document_text += cell.text + " "
                                document_text += "\n"
                    except ImportError:
                        logger.error("Библиотека python-docx не установлена")
                        return None
                    except Exception as e:
                        logger.error(f"Ошибка при обработке Word документа: {e}", exc_info=True)
                        return None
                
                # Обработка Excel файлов (.xlsx)
                elif file_ext_lower == '.xlsx':
                    try:
                        from openpyxl import load_workbook
                        document_file.seek(0)
                        workbook = load_workbook(document_file, data_only=True)
                        for sheet_name in workbook.sheetnames:
                            sheet = workbook[sheet_name]
                            document_text += f"\n--- Лист: {sheet_name} ---\n"
                            for row in sheet.iter_rows(values_only=True):
                                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                                if row_text.strip():
                                    document_text += row_text + "\n"
                    except ImportError:
                        logger.error("Библиотека openpyxl не установлена")
                        return None
                    except Exception as e:
                        logger.error(f"Ошибка при обработке Excel файла: {e}", exc_info=True)
                        return None
                
                # Обработка PowerPoint презентаций (.pptx)
                elif file_ext_lower == '.pptx':
                    try:
                        from pptx import Presentation
                        document_file.seek(0)
                        prs = Presentation(document_file)
                        for slide_num, slide in enumerate(prs.slides, 1):
                            document_text += f"\n--- Слайд {slide_num} ---\n"
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text.strip():
                                    document_text += shape.text + "\n"
                    except ImportError:
                        logger.error("Библиотека python-pptx не установлена")
                        return None
                    except Exception as e:
                        logger.error(f"Ошибка при обработке PowerPoint презентации: {e}", exc_info=True)
                        return None
                
                # Обработка OpenDocument Text (.odt)
                elif file_ext_lower == '.odt':
                    try:
                        from zipfile import ZipFile
                        import xml.etree.ElementTree as ET
                        document_file.seek(0)
                        with ZipFile(document_file, 'r') as odt_file:
                            content_xml = odt_file.read('content.xml')
                            root = ET.fromstring(content_xml)
                            # Простое извлечение текста из XML
                            for elem in root.iter():
                                if elem.text and elem.text.strip():
                                    document_text += elem.text.strip() + " "
                    except Exception as e:
                        logger.error(f"Ошибка при обработке ODT файла: {e}", exc_info=True)
                        return None
                
                else:
                    logger.warning(f"Неподдерживаемый формат документа: {file_extension}")
                    return None
                
                if not document_text.strip():
                    logger.warning(f"Из документа {file_extension} не удалось извлечь текст.")
                    return None
                
                # Ограничиваем размер текста (примерно 4000 токенов)
                if len(document_text) > 12000:
                    document_text = document_text[:12000] + "\n... (текст обрезан из-за ограничений)"
                
                logger.info(f"Текст из документа {file_extension} успешно извлечен ({len(document_text)} символов).")
                return document_text
                
    except asyncio.TimeoutError:
        logger.error(f"Таймаут при скачивании или обработке документа {file_extension} по URL: {document_url[:100]}...")
        return None
    except Exception as e:
        logger.error(f"Ошибка при обработке документа {file_extension} {document_url[:100]}...: {e}", exc_info=True)
        return None


async def analyze_pdf(pdf_text: str, openai_client) -> Optional[str]:
    """
    Анализирует PDF документ через OpenAI.

    :param pdf_text: Текст из PDF
    :param openai_client: Клиент OpenAI
    :return: Анализ документа или None при ошибке
    """
    try:
        logger.info(f"Начинаем анализ PDF текста ({len(pdf_text)} символов)...")
        # Оборачиваем синхронный вызов OpenAI в asyncio.to_thread для неблокирующего выполнения
        response = await asyncio.to_thread(
            openai_client.chat.completions.create,
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "user",
                    "content": f"""Проанализируй этот PDF документ и дай краткое описание на русском языке.

Включи в описание:
- Тип документа
- Основную тему/содержание
- Ключевые пункты
- Количество страниц (если видно из текста)

Текст документа:
{pdf_text}"""
                }
            ],
            max_tokens=500,
        )
        analysis = response.choices[0].message.content
        logger.info(f"Анализ PDF получен: {analysis[:100]}...")
        return analysis
    except Exception as e:
        logger.error(f"Ошибка при обращении к OpenAI для анализа PDF: {e}", exc_info=True)
        return None


async def transcribe_audio(bot: Bot, message: types.Message, openai_client) -> Optional[str]:
    """
    Транскрибирует голосовое или аудио сообщение через Whisper API.

    :param bot: Экземпляр бота
    :param message: Сообщение с аудио
    :param openai_client: Клиент OpenAI
    :return: Транскрибированный текст или None при ошибке
    """
    try:
        # Определяем file_id в зависимости от типа сообщения
        if message.voice:
            file_id = message.voice.file_id
            logger.info(f"Найдено голосовое сообщение, file_id: {file_id}")
        elif message.audio:
            file_id = message.audio.file_id
            logger.info(f"Найдено аудио сообщение, file_id: {file_id}")
        else:
            logger.warning("Сообщение не содержит голосового или аудио контента")
            return None

        file_info = await bot.get_file(file_id)
        file_url = f'https://api.telegram.org/file/bot{bot.token}/{file_info.file_path}'
        logger.info(f"Скачиваем аудио файл: {file_url}")

        # Скачиваем файл
        async with aiohttp.ClientSession() as session:
            async with session.get(file_url, timeout=aiohttp.ClientTimeout(total=60)) as response:
                response.raise_for_status()
                audio_content = await response.read()
                logger.info(f"Аудио файл скачан, размер: {len(audio_content)} байт")

        # Транскрибируем через OpenAI Whisper
        # OpenAI API требует файл в формате (filename, file_object) или (filename, file_object, content_type)
        audio_file = io.BytesIO(audio_content)
        filename = file_info.file_path.split('/')[-1] if file_info.file_path else "audio.ogg"
        
        # Убеждаемся, что файл имеет правильное расширение
        if not filename or not filename.lower().endswith(('.ogg', '.mp3', '.wav', '.m4a', '.mp4', '.flac', '.webm')):
            # Определяем расширение по типу сообщения
            if message.voice:
                filename = "voice.ogg"
            else:
                filename = "audio.ogg"
        
        logger.info(f"Отправляем файл в Whisper API: {filename}, размер: {len(audio_content)} байт")
        
        # Перематываем файл в начало
        audio_file.seek(0)
        
        # OpenAI SDK принимает файл как tuple (filename, file_object)
        # SDK автоматически определит content_type по расширению файла
        try:
            transcription = openai_client.audio.transcriptions.create(
                model="whisper-1",
                file=(filename, audio_file)
            )
        except Exception as api_error:
            logger.error(f"Ошибка при вызове Whisper API: {api_error}")
            # Пробуем без указания имени файла
            audio_file.seek(0)
            transcription = openai_client.audio.transcriptions.create(
                model="whisper-1",
                file=audio_file
            )

        transcribed_text = transcription.text
        logger.info(f"Транскрибация успешна: {transcribed_text[:100]}...")
        return transcribed_text
    except Exception as e:
        logger.error(f"Ошибка при транскрибации аудио: {e}", exc_info=True)
        return None


async def prepare_message_content(
    bot: Bot,
    message: types.Message,
    openai_client
) -> str:
    """
    Подготавливает полный контент сообщения для обработки AI.

    Обрабатывает различные типы контента: текст, фото, PDF, документы (txt, docx, xlsx, pptx, odt), 
    аудио, голосовые сообщения.

    :param bot: Экземпляр бота
    :param message: Сообщение для обработки
    :param openai_client: Клиент OpenAI
    :return: Полный текст для обработки AI
    """
    logger.info(f"Начинаем подготовку контента сообщения {message.message_id}")
    content_parts = []

    # Базовый текст (текст сообщения или подпись) - ВСЕГДА обрабатываем первым
    base_text = message.text or message.caption or ""
    if base_text:
        logger.info(f"Найден базовый текст: {base_text[:100]}...")
        try:
            # Обрабатываем URL в тексте (с таймаутом, чтобы не блокировать)
            base_text = await asyncio.wait_for(
                process_url_in_text(base_text),
                timeout=10.0
            )
            content_parts.append(base_text)
        except asyncio.TimeoutError:
            logger.warning("Таймаут при обработке URL в тексте, используем исходный текст")
            content_parts.append(base_text)  # Используем исходный текст без URL
        except Exception as e:
            logger.error(f"Ошибка при обработке URL в тексте: {e}, используем исходный текст")
            content_parts.append(base_text)  # Используем исходный текст без URL
    else:
        logger.info("Базовый текст отсутствует")

    # Обработка фотографий (неблокирующая, с таймаутом)
    if message.photo:
        logger.info("Обнаружено фото в сообщении, начинаем обработку")
        try:
            image_url = await asyncio.wait_for(get_photo_url(bot, message), timeout=10.0)
            if image_url:
                logger.info(f"URL изображения получен: {image_url[:100]}...")
                description = await asyncio.wait_for(
                    get_image_description(image_url, openai_client), 
                    timeout=30.0
                )
                if description:
                    logger.info(f"Описание изображения получено: {description[:100]}...")
                    content_parts.append(f"\nОписание изображения: {description}")
                else:
                    logger.warning("Не удалось получить описание изображения")
            else:
                logger.warning("Не удалось получить URL изображения")
        except asyncio.TimeoutError:
            logger.warning("Таймаут при обработке фото, пропускаем описание")
        except Exception as e:
            logger.error(f"Ошибка при обработке фото: {e}", exc_info=True)

    # Обработка PDF документов (неблокирующая, с таймаутом)
    # Извлекаем только текст из PDF, без анализа через OpenAI
    if message.document and message.document.file_name and message.document.file_name.lower().endswith('.pdf'):
        logger.info("Обнаружен PDF документ в сообщении, начинаем извлечение текста")
        try:
            file_info = await asyncio.wait_for(bot.get_file(message.document.file_id), timeout=10.0)
            pdf_url = f'https://api.telegram.org/file/bot{bot.token}/{file_info.file_path}'
            pdf_text = await asyncio.wait_for(extract_pdf_text(pdf_url), timeout=60.0)
            if pdf_text:
                logger.info(f"Текст из PDF извлечен: {len(pdf_text)} символов")
                content_parts.append(f"\n\nТекст из PDF документа:\n{pdf_text}")
            else:
                logger.warning("Не удалось извлечь текст из PDF документа")
        except asyncio.TimeoutError:
            logger.warning("Таймаут при обработке PDF, пропускаем извлечение текста")
        except Exception as e:
            logger.error(f"Ошибка при обработке PDF: {e}", exc_info=True)

    # Обработка других документов (txt, docx, xlsx, pptx, odt) - неблокирующая, с таймаутом
    # Извлекаем только текст из документов, без анализа через OpenAI
    if message.document and message.document.file_name:
        file_name_lower = message.document.file_name.lower()
        # Поддерживаемые форматы документов
        supported_extensions = ['.txt', '.docx', '.xlsx', '.pptx', '.odt']
        file_extension = None
        for ext in supported_extensions:
            if file_name_lower.endswith(ext):
                file_extension = ext
                break
        
        if file_extension:
            logger.info(f"Обнаружен документ {file_extension} в сообщении, начинаем извлечение текста")
            try:
                file_info = await asyncio.wait_for(bot.get_file(message.document.file_id), timeout=10.0)
                document_url = f'https://api.telegram.org/file/bot{bot.token}/{file_info.file_path}'
                document_text = await asyncio.wait_for(
                    extract_document_text(document_url, file_extension), 
                    timeout=60.0
                )
                if document_text:
                    logger.info(f"Текст из документа {file_extension} извлечен: {len(document_text)} символов")
                    content_parts.append(f"\n\nТекст из документа {file_extension}:\n{document_text}")
                else:
                    logger.warning(f"Не удалось извлечь текст из документа {file_extension}")
            except asyncio.TimeoutError:
                logger.warning(f"Таймаут при обработке документа {file_extension}, пропускаем извлечение текста")
            except Exception as e:
                logger.error(f"Ошибка при обработке документа {file_extension}: {e}", exc_info=True)

    # ОТКЛЮЧЕНО ДЛЯ ОТЛАДКИ: Обработка голосовых сообщений
    # if message.voice:
    #     logger.info("Обнаружено голосовое сообщение, начинаем транскрибацию")
    #     try:
    #         transcription = await asyncio.wait_for(
    #             transcribe_audio(bot, message, openai_client), 
    #             timeout=60.0
    #         )
    #         if transcription:
    #             logger.info(f"Транскрибация голосового сообщения успешна: {transcription[:100]}...")
    #             content_parts.append(f"\nТранскрипция аудио: {transcription}")
    #         else:
    #             logger.warning("Не удалось транскрибировать голосовое сообщение")
    #     except asyncio.TimeoutError:
    #         logger.warning("Таймаут при транскрибации голосового сообщения, пропускаем")
    #     except Exception as e:
    #         logger.error(f"Ошибка при транскрибации голосового сообщения: {e}", exc_info=True)
    if message.voice:
        logger.info("ОТЛАДКА: Транскрибация голосовых сообщений отключена")

    # ОТКЛЮЧЕНО ДЛЯ ОТЛАДКИ: Обработка аудио файлов
    # if message.audio:
    #     logger.info("Обнаружено аудио сообщение, начинаем транскрибацию")
    #     try:
    #         transcription = await asyncio.wait_for(
    #             transcribe_audio(bot, message, openai_client), 
    #             timeout=60.0
    #         )
    #         if transcription:
    #             logger.info(f"Транскрибация аудио файла успешна: {transcription[:100]}...")
    #             content_parts.append(f"\nТранскрипция аудио: {transcription}")
    #         else:
    #             logger.warning("Не удалось транскрибировать аудио файл")
    #     except asyncio.TimeoutError:
    #         logger.warning("Таймаут при транскрибации аудио файла, пропускаем")
    #     except Exception as e:
    #         logger.error(f"Ошибка при транскрибации аудио файла: {e}", exc_info=True)
    if message.audio:
        logger.info("ОТЛАДКА: Транскрибация аудио файлов отключена")

    # Обработка опросов
    if message.poll:
        logger.info("Обнаружен опрос в сообщении")
        content_parts.append(f"Опрос: {message.poll.question}")

    # Если нет контента, но есть базовый текст или медиа, возвращаем хотя бы описание типа контента
    if not content_parts:
        if message.photo:
            result = "Фото без подписи"
        elif message.voice:
            result = "Голосовое сообщение"
        elif message.audio:
            result = "Аудио файл"
        elif message.video:
            result = message.caption or "Видео без подписи"
        elif message.document:
            result = message.caption or f"Документ: {message.document.file_name or 'без имени'}"
        else:
            result = "Сообщение без текстового контента"
        logger.info(f"Подготовка контента завершена. Результат (только тип контента): {result}")
    else:
        result = "\n".join(content_parts)
        logger.info(f"Подготовка контента завершена. Результат ({len(result)} символов): {result[:200]}...")
    
    return result
